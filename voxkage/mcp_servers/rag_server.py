"""
MCP Server: VoxKage Local Knowledge RAG (voxkage-rag)

Deep document memory powered by SQLite FTS5.
Indexes PDFs, Word docs, Excel, PowerPoint, TXT, and code files.
Detects file changes via SHA256 hash and auto-reindexes on demand.

Tools:
  index_document(file_path)          — Index or re-index a file into the RAG
  query_rag(query, top_k, file_filter) — Keyword search across all indexed documents
  check_and_index(file_path)         — Auto-index if new/changed, skip if unchanged
  list_indexed_documents()           — Show all indexed files and their metadata
  delete_from_rag(file_path)         — Remove a file from the index
  index_directory(directory, extensions) — Bulk index an entire folder/codebase

Run standalone: python mcp_servers/rag_server.py
"""

import os
import sys
import hashlib
import json
import re
import time
import logging
import sqlite3
from pathlib import Path
from datetime import datetime
import contextlib

@contextlib.contextmanager
def suppress_stdout():
    """Redirect stdout to /dev/null — prevents progress bars from corrupting MCP stdio JSON-RPC."""
    with open(os.devnull, "w") as devnull:
        old_stdout = sys.stdout
        sys.stdout = devnull
        try:
            yield
        finally:
            sys.stdout = old_stdout

@contextlib.contextmanager
def suppress_all_output():
    """Redirect both stdout and stderr — used during heavy ops to keep MCP channel clean."""
    with open(os.devnull, "w") as devnull:
        old_stdout, old_stderr = sys.stdout, sys.stderr
        sys.stdout = devnull
        sys.stderr = devnull
        try:
            yield
        finally:
            sys.stdout = old_stdout
            sys.stderr = old_stderr

_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)

from voxkage._env import load_voxkage_env
load_voxkage_env()

from mcp.server.fastmcp import FastMCP
mcp = FastMCP("voxkage-rag")
logger = logging.getLogger(__name__)

from voxkage.paths import rag_dir, brain_dir

# ── Storage Paths ─────────────────────────────────────────────────────────────
_RAG_DIR   = str(rag_dir())
_DB_FILE   = os.path.join(_RAG_DIR, "rag_database.db")
_META_FILE = os.path.join(_RAG_DIR, "index_meta.json")  # Keep reference if needed
os.makedirs(_RAG_DIR, exist_ok=True)

# ── Supported file extensions ─────────────────────────────────────────────────
_TEXT_EXTS = {
    ".txt", ".md", ".csv", ".log", ".yaml", ".yml", ".json",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".html", ".htm",
    ".java", ".cpp", ".c", ".h", ".cs", ".go", ".rs", ".rb",
    ".php", ".sh", ".bat", ".ps1", ".sql", ".xml", ".toml", ".ini",
    ".gitignore", ".env",
}
_RICH_EXTS = {".pdf", ".docx", ".xlsx", ".pptx", ".xls", ".doc", ".png", ".jpg", ".jpeg", ".bmp"}
_ALL_EXTS  = _TEXT_EXTS | _RICH_EXTS

# Skip these in bulk directory indexing
_SKIP_DIRS = {
    "__pycache__", ".git", "node_modules", ".venv", "venv", "env",
    "dist", "build", ".next", ".nuxt", ".svelte-kit", "target",
    ".idea", ".vscode", ".terraform", ".serverless", ".expo",
    "vendor", ".cache", ".gemini", ".antigravitycli", "out",
}

_SKIP_FILES = {
    "package-lock.json", "yarn.lock", "pnpm-lock.yaml", "cargo.lock",
    "gemfile.lock", "composer.lock", "poetry.lock", "mix.lock",
    "pipfile.lock", "tsconfig.tsbuildinfo", "gradle-wrapper.properties",
}


# ═════════════════════════════════════════════════════════════════════════════
# INTERNAL HELPERS & DB SETUP
# ═════════════════════════════════════════════════════════════════════════════

@contextlib.contextmanager
def _get_db():
    """Context manager for SQLite database connection."""
    conn = sqlite3.connect(_DB_FILE, timeout=30.0)
    conn.row_factory = sqlite3.Row
    try:
        yield conn
    finally:
        conn.close()


def _init_db():
    """Initialize SQLite database with WAL mode and necessary tables."""
    os.makedirs(_RAG_DIR, exist_ok=True)
    with sqlite3.connect(_DB_FILE) as conn:
        conn.execute("PRAGMA journal_mode=WAL;")
        conn.execute("PRAGMA synchronous=NORMAL;")
        
        # Create metadata table
        conn.execute("""
            CREATE TABLE IF NOT EXISTS file_meta (
                file_path TEXT PRIMARY KEY,
                hash TEXT NOT NULL,
                indexed_at TEXT NOT NULL,
                chunks_count INTEGER NOT NULL,
                file_name TEXT NOT NULL,
                file_ext TEXT NOT NULL,
                file_size_kb REAL NOT NULL
            )
        """)
        
        # Create chunks FTS5 virtual table
        conn.execute("""
            CREATE VIRTUAL TABLE IF NOT EXISTS chunks USING fts5(
                content,
                file_path UNINDEXED,
                file_name UNINDEXED,
                chunk_index UNINDEXED
            )
        """)


# Initialize DB on import
_init_db()


def _sanitize_fts_query(query: str) -> str:
    """
    Sanitize natural language queries for SQLite FTS5.
    Extracts words (treating underscores and dots as separators) and formats them with prefix wildcards joined by OR.
    """
    query_clean = query.replace("_", " ").replace(".", " ")
    words = re.findall(r'\b\w+\b', query_clean)
    if not words:
        return ""
    return " OR ".join(f"{w}*" for w in words)


def _get_ocr():
    """Return RapidOCR instance (lazy-initialized singleton)."""
    if not hasattr(_get_ocr, "_instance"):
        try:
            from rapidocr_onnxruntime import RapidOCR
            _get_ocr._instance = RapidOCR()
        except ImportError:
            _get_ocr._instance = None
    return _get_ocr._instance


def _sha256(file_path: str) -> str:
    """Compute SHA256 hash of a file's content."""
    h = hashlib.sha256()
    try:
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return ""


def _load_meta() -> dict:
    """Load the metadata registry {file_path -> {hash, indexed_at, chunks}} from DB."""
    try:
        meta = {}
        with _get_db() as conn:
            rows = conn.execute("""
                SELECT file_path, hash, indexed_at, chunks_count, file_name, file_ext, file_size_kb
                FROM file_meta
            """).fetchall()
            for r in rows:
                meta[r["file_path"]] = {
                    "hash": r["hash"],
                    "indexed_at": r["indexed_at"],
                    "chunks": r["chunks_count"],
                    "file_name": r["file_name"],
                    "file_ext": r["file_ext"],
                    "file_size_kb": r["file_size_kb"]
                }
        return meta
    except Exception as e:
        logger.warning(f"Failed to load RAG meta: {e}")
        return {}


def _save_meta(meta: dict):
    """Save metadata registry dict to DB (for backward compatibility)."""
    try:
        with _get_db() as conn:
            conn.execute("BEGIN TRANSACTION;")
            try:
                for file_path, info in meta.items():
                    conn.execute("""
                        INSERT OR REPLACE INTO file_meta (file_path, hash, indexed_at, chunks_count, file_name, file_ext, file_size_kb)
                        VALUES (?, ?, ?, ?, ?, ?, ?)
                    """, (
                        file_path,
                        info.get("hash", ""),
                        info.get("indexed_at", ""),
                        info.get("chunks", 0),
                        info.get("file_name", os.path.basename(file_path)),
                        info.get("file_ext", Path(file_path).suffix),
                        info.get("file_size_kb", 0.0)
                    ))
                conn.commit()
            except Exception as e:
                conn.rollback()
                raise e
    except Exception as e:
        logger.warning(f"Failed to save RAG meta: {e}")


def _extract_text(file_path: str) -> str:
    """Extract full text from a file based on its extension."""
    ext = Path(file_path).suffix.lower()

    try:
        if ext in _TEXT_EXTS:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        elif ext == ".pdf":
            try:
                import fitz
            except ImportError:
                return "[RAG] PDF reading requires PyMuPDF. Run: pip install voxkage[docs_plus]"
            text = ""
            with fitz.open(file_path) as doc:
                for i, page in enumerate(doc):
                    page_text = page.get_text().strip()
                    # Only do OCR if strictly no text is found, to avoid hanging
                    if len(page_text) < 10:
                        try:
                            import numpy as np
                            ocr = _get_ocr()
                            if ocr:
                                with suppress_all_output():
                                    pix = page.get_pixmap(dpi=150, alpha=False)
                                    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
                                    result, _ = ocr(img)
                                if result:
                                    page_text = "\n".join([line[1] for line in result])
                        except Exception as e:
                            logger.error(f"OCR failed on PDF page {i+1}: {e}")
                    text += page_text + "\n"
            return text.strip()

        elif ext in (".png", ".jpg", ".jpeg", ".bmp"):
            try:
                import cv2
                ocr = _get_ocr()
                if not ocr:
                    return ""
                img = cv2.imread(file_path)
                result, _ = ocr(img)
                if result:
                    return "\n".join([line[1] for line in result])
                return ""
            except Exception as e:
                logger.error(f"OCR failed on image {file_path}: {e}")
                return ""

        elif ext in (".docx", ".doc"):
            from docx import Document
            doc = Document(file_path)
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)
            return "\n".join(parts)

        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                parts = []
                for sheet in wb.worksheets:
                    parts.append(f"[Sheet: {sheet.title}]")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join(
                            str(c) for c in row if c is not None and str(c).strip()
                        )
                        if row_text:
                            parts.append(row_text)
                return "\n".join(parts)
            except Exception:
                return f"Could not extract Excel content from {file_path}"

        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"[Slide {i}]")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text.strip())
                return "\n".join(parts)
            except Exception:
                return f"Could not extract PowerPoint content from {file_path}"

        else:
            return f"Unsupported file type: {ext}"

    except Exception as e:
        return f"Error extracting text: {e}"


def _smart_chunk(text: str, file_path: str, chunk_size: int = 800, overlap: int = 150) -> list[str]:
    """
    Split text into overlapping chunks.
    For code files: split on function/class boundaries when possible.
    For prose: split on paragraph boundaries.
    """
    ext = Path(file_path).suffix.lower()

    if not text or not text.strip():
        return []

    if ext in {".py", ".js", ".ts", ".java", ".cpp", ".c", ".cs", ".go", ".rs", ".rb", ".php"}:
        pattern = r"(?m)(?=^(?:def |class |function |async def |export function |func |fn ))"
        logical_blocks = re.split(pattern, text)
        chunks = []
        current = ""
        for block in logical_blocks:
            if len(current) + len(block) > chunk_size:
                if current.strip():
                    chunks.append(current.strip())
                current = block
            else:
                current += block
        if current.strip():
            chunks.append(current.strip())
        if chunks:
            return chunks

    paragraphs = re.split(r"\n{2,}", text)
    chunks = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) > chunk_size:
            if current.strip():
                chunks.append(current.strip())
            current = current[-overlap:] + "\n\n" + para if current else para
        else:
            current += ("\n\n" if current else "") + para
    if current.strip():
        chunks.append(current.strip())

    if not chunks:
        chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size - overlap)]

    return [c for c in chunks if len(c.strip()) > 30]


def _safe_id(file_path: str, chunk_idx: int) -> str:
    """Generate a stable, ChromaDB-safe document ID (kept for signature compatibility)."""
    key = f"{file_path}::chunk_{chunk_idx}"
    return hashlib.md5(key.encode()).hexdigest()


def _do_index(file_path: str, force: bool = False) -> dict:
    """
    Core indexing logic. Returns a status dict.
    If force=True, always re-indexes even if hash matches.
    """
    file_path = os.path.normpath(os.path.abspath(file_path))

    if not os.path.exists(file_path):
        return {"status": "error", "message": f"File not found: {file_path}"}

    ext = Path(file_path).suffix.lower()
    if ext not in _ALL_EXTS:
        return {"status": "skipped", "message": f"Unsupported extension: {ext}"}

    # Size check before text extraction to avoid CPU/memory hangs
    file_size = os.path.getsize(file_path)
    if ext in _TEXT_EXTS:
        max_size = 512 * 1024  # 500 KB limit for code/text
        if file_size > max_size:
            return {
                "status": "skipped",
                "message": f"Text/code file size ({round(file_size/1024, 1)} KB) exceeds 500 KB limit."
            }
    elif ext in _RICH_EXTS:
        max_size = 5 * 1024 * 1024  # 5 MB limit for rich docs/media
        if file_size > max_size:
            return {
                "status": "skipped",
                "message": f"Rich document size ({round(file_size/(1024*1024), 1)} MB) exceeds 5 MB limit."
            }

    current_hash = _sha256(file_path)
    
    # Query file_meta for hash check
    existing = None
    try:
        with _get_db() as conn:
            row = conn.execute(
                "SELECT hash, indexed_at, chunks_count, file_name FROM file_meta WHERE file_path = ?",
                (file_path,)
            ).fetchone()
            if row:
                existing = {
                    "hash": row["hash"],
                    "indexed_at": row["indexed_at"],
                    "chunks": row["chunks_count"],
                    "file_name": row["file_name"]
                }
    except Exception as e:
        logger.warning(f"Database error during hash check: {e}")

    if not force and existing and existing["hash"] == current_hash:
        return {
            "status": "unchanged",
            "message": f"File unchanged since {existing['indexed_at']}. Using cached index.",
            "chunks": existing["chunks"],
            "file": existing["file_name"],
        }

    text = _extract_text(file_path)
    if not text:
        return {"status": "error", "message": "No content extracted from file"}
    if text.startswith("Error extracting text:") or text.startswith("[RAG] PDF reading requires") or text.startswith("Unsupported file type:"):
        return {"status": "error", "message": text}

    chunks = _smart_chunk(text, file_path)
    if not chunks:
        return {"status": "error", "message": "No content extracted from file"}

    # Cap chunks to prevent index size explosion
    if len(chunks) > 250:
        chunks = chunks[:250]

    # Time tracking
    start_time = time.time()
    file_name = os.path.basename(file_path)
    now = datetime.now().isoformat()

    try:
        with _get_db() as conn:
            conn.execute("BEGIN TRANSACTION;")
            try:
                # Delete old chunks
                conn.execute("DELETE FROM chunks WHERE file_path = ?", (file_path,))
                
                # Insert new chunks
                for i, chunk in enumerate(chunks):
                    conn.execute(
                        "INSERT INTO chunks (content, file_path, file_name, chunk_index) VALUES (?, ?, ?, ?)",
                        (chunk, file_path, file_name, i)
                    )
                
                # Upsert file_meta
                conn.execute("""
                    INSERT OR REPLACE INTO file_meta (file_path, hash, indexed_at, chunks_count, file_name, file_ext, file_size_kb)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                """, (file_path, current_hash, now, len(chunks), file_name, ext, round(file_size / 1024, 1)))
                
                conn.commit()
            except Exception as e:
                conn.rollback()
                raise e
    except Exception as e:
        return {"status": "error", "message": f"Database indexing failed: {e}"}

    action = "reindexed" if existing else "indexed"
    duration = round(time.time() - start_time, 2)
    return {
        "status": "success",
        "action": action,
        "file": file_name,
        "chunks": len(chunks),
        "characters": len(text),
        "message": f"[OK] {action.capitalize()} '{file_name}' into RAG ({len(chunks)} chunks from {len(text):,} chars) in {duration}s",
    }


# ═════════════════════════════════════════════════════════════════════════════
# MCP TOOLS
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def index_document(file_path: str) -> str:
    """
    RAG MEMORY: Index or re-index a document into VoxKage's local knowledge base.

    Automatically detects if the file is NEW (indexes fresh) or CHANGED (reindexes)
    or UNCHANGED (returns cached result instantly — no work done).

    Supports: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), TXT, CSV,
              Python, JavaScript, TypeScript, Java, C/C++, Go, Rust, SQL, Markdown, YAML, JSON, and more.

    ALWAYS call this after reading any file for the first time so VoxKage remembers it.
    """
    result = _do_index(file_path, force=False)
    return result.get("message", str(result))


@mcp.tool()
def check_and_index(file_path: str) -> str:
    """
    RAG MEMORY: Smart auto-index gate. Call this BEFORE reading ANY file.

    - If file is new → indexes it automatically
    - If file has changed since last index → reindexes it
    - If file is unchanged → returns instantly from cache

    This is the STANDARD workflow for file access:
    1. check_and_index(file_path)   ← always call this first
    2. query_rag(your_question)     ← search the indexed content
    """
    result = _do_index(file_path, force=False)
    status = result.get("status", "error")
    if status == "unchanged":
        return f"[RAG] Already indexed and unchanged — using cache. {result['chunks']} chunks available for '{result['file']}'."
    elif status == "success":
        return f"[RAG] {result['message']}"
    elif status == "skipped":
        return f"[RAG SKIP] {result['message']} — will read file directly."
    else:
        return f"[RAG ERROR] {result.get('message', 'Unknown error')}"


@mcp.tool()
def query_rag(
    query: str,
    top_k: int = 5,
    file_filter: str = "",
) -> str:
    """
    RAG MEMORY: Keyword search across all indexed documents using SQLite FTS5.

    Use this to answer questions about ANY file VoxKage has seen before,
    without needing to re-read the file.

    Parameters:
      query       : Natural language question or keyword search
      top_k       : Number of results to return (default 5)
      file_filter : Optional — restrict search to files matching this name/path substring
    """
    query_sanitized = _sanitize_fts_query(query)
    if not query_sanitized:
        return "[RAG] Query contains no valid search terms."

    try:
        with _get_db() as conn:
            # Check if any documents are indexed
            total_chunks = conn.execute("SELECT COUNT(*) FROM chunks").fetchone()[0]
            if total_chunks == 0:
                return "[RAG] No documents indexed yet. Use index_document() or check_and_index() first."

            # Generate dynamic summing filename boost
            query_clean = query.replace("_", " ").replace(".", " ")
            words = re.findall(r'\b\w+\b', query_clean)
            boost_sql = []
            boost_params = []
            for w in words:
                if len(w) > 2:
                    boost_sql.append("CASE WHEN c.file_name LIKE ? THEN 50.0 ELSE 0.0 END")
                    boost_params.append(f"%{w}%")

            boost_clause = ""
            if boost_sql:
                boost_clause = f" - ({' + '.join(boost_sql)})"

            # Construct SQL query
            sql = f"""
                SELECT c.content, c.file_path, c.file_name, c.chunk_index, 
                       (c.rank{boost_clause}) as rank, f.chunks_count
                FROM chunks c
                JOIN file_meta f ON c.file_path = f.file_path
                WHERE chunks MATCH ?
            """
            params = boost_params + [query_sanitized]

            if file_filter:
                sql += " AND (f.file_path LIKE ? OR f.file_name LIKE ?)"
                params.extend([f"%{file_filter}%", f"%{file_filter}%"])

            sql += " ORDER BY rank LIMIT ?"
            params.append(top_k)

            rows = conn.execute(sql, params).fetchall()

    except sqlite3.OperationalError:
        # Fallback query if MATCH query syntax fails
        try:
            query_clean = query.replace("_", " ").replace(".", " ")
            words = re.findall(r'\b\w+\b', query_clean)
            simple_query = " OR ".join(words)
            
            boost_sql = []
            boost_params = []
            for w in words:
                if len(w) > 2:
                    boost_sql.append("CASE WHEN c.file_name LIKE ? THEN 50.0 ELSE 0.0 END")
                    boost_params.append(f"%{w}%")

            boost_clause = ""
            if boost_sql:
                boost_clause = f" - ({' + '.join(boost_sql)})"

            with _get_db() as conn:
                sql = f"""
                    SELECT c.content, c.file_path, c.file_name, c.chunk_index, 
                           (c.rank{boost_clause}) as rank, f.chunks_count
                    FROM chunks c
                    JOIN file_meta f ON c.file_path = f.file_path
                    WHERE chunks MATCH ?
                """
                params = boost_params + [simple_query]

                if file_filter:
                    sql += " AND (f.file_path LIKE ? OR f.file_name LIKE ?)"
                    params.extend([f"%{file_filter}%", f"%{file_filter}%"])

                sql += " ORDER BY rank LIMIT ?"
                params.append(top_k)

                rows = conn.execute(sql, params).fetchall()
        except Exception as e2:
            return f"[RAG ERROR] Query failed: {e2}"
    except Exception as e:
        return f"[RAG ERROR] Query failed: {e}"

    if not rows:
        return "[RAG] No relevant results found for that query."

    lines = [f"[RAG] Found {len(rows)} relevant chunk(s) for: \"{query}\"\n"]
    for i, row in enumerate(rows, 1):
        score = round(-row["rank"], 2)
        fn = row["file_name"]
        chunk_n = row["chunk_index"]
        total = row["chunks_count"]
        lines.append(f"--- Result {i} | {fn} (chunk {chunk_n + 1}/{total}) | Match Score: {score} ---")
        lines.append(row["content"].strip())
        lines.append("")

    return "\n".join(lines)


@mcp.tool()
def list_indexed_documents() -> str:
    """
    RAG MEMORY: List all documents currently in VoxKage's knowledge base.

    Shows: file name, size, chunk count, when indexed, and whether the file has
    changed on disk since it was last indexed (needs reindex if so).
    """
    try:
        with _get_db() as conn:
            rows = conn.execute("""
                SELECT file_path, hash, indexed_at, chunks_count, file_name, file_ext, file_size_kb
                FROM file_meta
                ORDER BY indexed_at DESC
            """).fetchall()
    except Exception as e:
        return f"[RAG ERROR] Failed to list documents: {e}"

    if not rows:
        return "[RAG] No documents indexed yet. Use index_document() to start building your knowledge base."

    lines = [f"[RAG] VoxKage Knowledge Base — {len(rows)} indexed document(s)\n"]
    lines.append(f"{'File':<45} {'Chunks':>6}  {'Size KB':>8}  {'Indexed At':<22}  {'Status'}")
    lines.append("-" * 110)

    for row in rows:
        file_path = row["file_path"]
        fn = row["file_name"][:44]
        chunks = row["chunks_count"]
        size_kb = row["file_size_kb"]
        indexed_at = row["indexed_at"][:19]

        status = "[OK] current"
        if not os.path.exists(file_path):
            status = "[!!] FILE MISSING"
        elif _sha256(file_path) != row["hash"]:
            status = "[!!] CHANGED -- needs reindex"

        lines.append(f"{fn:<45} {str(chunks):>6}  {str(size_kb):>8}  {indexed_at:<22}  {status}")

    return "\n".join(lines)


@mcp.tool()
def delete_from_rag(file_path: str) -> str:
    """
    RAG MEMORY: Remove a document from VoxKage's knowledge base.
    Use this if a file has been deleted or you no longer want it indexed.
    """
    file_path = os.path.normpath(os.path.abspath(file_path))
    
    with _get_db() as conn:
        row = conn.execute(
            "SELECT file_name, chunks_count FROM file_meta WHERE file_path = ?", 
            (file_path,)
        ).fetchone()
        if not row:
            return f"[RAG] '{os.path.basename(file_path)}' is not in the index."

        file_name = row["file_name"]
        chunk_count = row["chunks_count"]

        conn.execute("BEGIN TRANSACTION;")
        try:
            conn.execute("DELETE FROM chunks WHERE file_path = ?", (file_path,))
            conn.execute("DELETE FROM file_meta WHERE file_path = ?", (file_path,))
            conn.commit()
        except Exception as e:
            conn.rollback()
            return f"[RAG ERROR] Failed to delete '{file_name}': {e}"

    return f"[RAG] Removed '{file_name}' and its {chunk_count} chunks from the knowledge base."


@mcp.tool()
def index_directory(
    directory: str,
    extensions: str = "",
    recursive: bool = True,
) -> str:
    """
    RAG MEMORY: Bulk-index an entire directory (codebase, documents folder, etc.).

    Parameters:
      directory  : The folder to index (e.g. "C:\\Users\\<Username>\\Desktop\\Vision-Assistant")
      extensions : Comma-separated list of file extensions to include (e.g. ".py,.md,.txt")
                   Leave empty to use all supported extensions.
      recursive  : Whether to index subdirectories (default True)

    This is how you index an ENTIRE CODEBASE or documents folder.
    Already-indexed unchanged files are skipped automatically (only new/changed files are processed).
    """
    directory = os.path.normpath(os.path.abspath(directory))
    if not os.path.isdir(directory):
        return f"[RAG ERROR] Not a directory: {directory}"

    if extensions and isinstance(extensions, str) and extensions.strip():
        target_exts = {e.strip().lower() if e.strip().startswith(".") else f".{e.strip().lower()}"
                       for e in extensions.split(",") if e.strip()}
        target_exts = target_exts & _ALL_EXTS
    else:
        target_exts = _TEXT_EXTS

    all_target_files = []
    walker = os.walk(directory) if recursive else [(directory, [], os.listdir(directory))]
    for root, dirs, files in walker:
        if recursive:
            dirs[:] = [d for d in dirs if d not in _SKIP_DIRS and not d.startswith(".")]
        for fname in files:
            if fname.lower() in _SKIP_FILES or fname.startswith("."):
                continue
            ext = Path(fname).suffix.lower()
            if ext in target_exts:
                all_target_files.append(os.path.join(root, fname))
    
    total_files = len(all_target_files)
    if total_files == 0:
        return f"[RAG] No supported files found in: {directory}\n  Looked for extensions: {', '.join(sorted(target_exts))}"

    results = {"indexed": 0, "reindexed": 0, "unchanged": 0, "skipped": 0, "errors": 0}
    processed = []

    for idx, fpath in enumerate(all_target_files, 1):
        fname = os.path.basename(fpath)
        try:
            result = _do_index(fpath, force=False)
            status = result.get("status", "error")
            action = result.get("action", "")
            if status == "success" and action == "indexed":
                results["indexed"] += 1
                processed.append(f"  [NEW]: {fname}")
            elif status == "success" and action == "reindexed":
                results["reindexed"] += 1
                processed.append(f"  [UPDATED]: {fname}")
            elif status == "unchanged":
                results["unchanged"] += 1
            elif status == "skipped":
                results["skipped"] += 1
            else:
                results["errors"] += 1
                processed.append(f"  [ERROR]: {fname} -- {result.get('message', '')}")
        except Exception as e:
            results["errors"] += 1
            processed.append(f"  [EXCEPTION]: {fname} -- {e}")

    summary_lines = [
        f"[RAG] ✓ Directory indexing complete for: {directory}",
        f"  Total files scanned : {total_files}",
        f"  New files indexed  : {results['indexed']}",
        f"  Files reindexed  : {results['reindexed']}",
        f"  Unchanged (cached): {results['unchanged']}",
        f"  Unsupported/skipped: {results['skipped']}",
        f"  Errors           : {results['errors']}",
        "",
    ]
    if processed:
        summary_lines.append("Changes:")
        summary_lines.extend(processed[:30])
        if len(processed) > 30:
            summary_lines.append(f"  ... and {len(processed) - 30} more")

    try:
        _BRAIN_DIR = brain_dir()
        _cache_file = _BRAIN_DIR / "index_cache.json"
        import json as _json
        _cache = {}
        if _cache_file.exists():
            try:
                _cache = _json.loads(_cache_file.read_text(encoding="utf-8"))
            except Exception:
                pass
        _cache[os.path.normpath(directory)] = time.time()
        _cache_file.write_text(_json.dumps(_cache, indent=2), encoding="utf-8")
    except Exception:
        pass

    return "\n".join(summary_lines)


if __name__ == "__main__":
    mcp.run()
