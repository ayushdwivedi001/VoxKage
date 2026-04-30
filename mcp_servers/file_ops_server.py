"""
MCP Server: File Operations — Create, Edit, Delete, Convert

Supports: Word (.docx), Excel (.xlsx), PPT (.pptx), TXT, CSV, PDF,
          HTML, JSON, JS/JSX/TS/TSX, Python, YAML, CSS, Markdown,
          and any plain-text format. Also creates folders.

All destructive operations require confirmed=True after user says "yes/agreed".

Run standalone: python mcp_servers/file_ops_server.py
"""

import os
import re
import sys

_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)

from dotenv import load_dotenv
load_dotenv(os.path.join(_ROOT, ".env"))

from mcp.server.fastmcp import FastMCP
mcp = FastMCP("voxkage-fileops")


# ── Directory resolution (self-contained) ────────────────────────────────────

_NOISE = {"the","a","an","my","is","in","on","at","file","folder","app","application",
          "of","to","open","go","find","show","me","please","can","you","it","i",
          "want","that","with","inside","into","under","new","create","make"}

def _tokens(text):
    return set(re.sub(r"[^a-z0-9 ]", " ", text.lower()).split()) - _NOISE

def _fscore(query, name):
    q, n = _tokens(query), _tokens(name)
    if not q: return 0.0
    if query.lower() in name.lower() or name.lower() in query.lower(): return 1.0
    if not n: return 0.0
    return len(q & n) / len(q)


def _resolve_dir(description: str) -> str | None:
    """Find a real directory path from a natural language description."""
    from config_loader import load_config
    config = load_config()
    user_home = os.path.expanduser("~")
    T = 0.5

    # 1. Absolute path given directly
    if os.path.isdir(description):
        return description

    # 2. Config app_launch_commands paths
    for alias, cmd in config.get("app_launch_commands", {}).items():
        if _fscore(description, alias) >= T:
            parts = cmd.split('"')
            for p in parts:
                if len(p) > 3 and (":" in p or "\\" in p):
                    candidate = p.strip()
                    if os.path.isdir(candidate):
                        return candidate

    # 2. Common user folders (checked BEFORE C:\ root to avoid false matches)
    for name in ("Desktop", "Documents", "Downloads", "Pictures", "Videos", "Music"):
        path = os.path.join(user_home, name)
        if _fscore(description, name) >= T and os.path.isdir(path):
            return path

    # 3. C:\ root folders
    try:
        for entry in os.scandir("C:\\"):
            if entry.is_dir() and _fscore(description, entry.name) >= T:
                return entry.path
    except Exception:
        pass

    return None


def _find_file_in_dir(directory: str, keyword: str) -> str | None:
    """Find a file by keyword inside a specific directory (non-recursive)."""
    try:
        for entry in os.scandir(directory):
            if entry.is_file() and _fscore(keyword, entry.name) >= 0.4:
                return entry.path
    except Exception:
        pass
    # Recursive one level
    try:
        for entry in os.scandir(directory):
            if entry.is_dir():
                for sub in os.scandir(entry.path):
                    if sub.is_file() and _fscore(keyword, sub.name) >= 0.4:
                        return sub.path
    except Exception:
        pass
    return None


# ── File type helpers ─────────────────────────────────────────────────────────

_EXT_MAP = {
    "word": ".docx", "excel": ".xlsx", "csv": ".csv",
    "txt": ".txt", "pdf": ".pdf", "text": ".txt",
    "ppt": ".pptx", "pptx": ".pptx", "powerpoint": ".pptx",
    "html": ".html", "json": ".json", "python": ".py",
    "javascript": ".js", "js": ".js", "jsx": ".jsx",
    "typescript": ".ts", "ts": ".ts", "tsx": ".tsx",
    "css": ".css", "yaml": ".yaml", "markdown": ".md", "md": ".md",
    "sql": ".sql", "xml": ".xml", "sh": ".sh", "bat": ".bat",
    "folder": "", "directory": "",
}

_TYPE_FROM_EXT = {
    ".docx": "word", ".doc": "word",
    ".xlsx": "excel", ".xls": "excel",
    ".csv": "csv", ".txt": "txt", ".pdf": "pdf",
    ".pptx": "pptx", ".ppt": "pptx",
    ".md": "txt", ".py": "txt", ".html": "txt", ".json": "txt",
    ".js": "txt", ".jsx": "txt", ".ts": "txt", ".tsx": "txt",
    ".css": "txt", ".yaml": "txt", ".yml": "txt",
    ".sql": "txt", ".xml": "txt", ".sh": "txt", ".bat": "txt",
}

_TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".html", ".json", ".js", ".jsx",
    ".ts", ".tsx", ".css", ".yaml", ".yml", ".sql", ".xml",
    ".sh", ".bat", ".env", ".toml", ".ini", ".cfg", ".log",
    ".cs", ".java", ".cpp", ".c", ".h", ".rb", ".go", ".rs",
    ".vue", ".svelte", ".jsx", ".tsx",
}

def _infer_type(filename: str, hint: str = "auto") -> str:
    if hint != "auto" and hint.lower() in _EXT_MAP:
        return hint.lower()
    ext = os.path.splitext(filename.lower())[1]
    return _TYPE_FROM_EXT.get(ext, "txt")

def _ensure_ext(filename: str, ftype: str) -> str:
    if "." not in os.path.basename(filename) and ftype not in ("folder", "directory"):
        return filename + _EXT_MAP.get(ftype, ".txt")
    return filename


# ── Word document helpers ─────────────────────────────────────────────────────

def _inline_fmt(paragraph, text: str):
    """Add text to paragraph with **bold** and *italic* support."""
    parts = re.split(r'\*\*(.+?)\*\*', text)
    for i, part in enumerate(parts):
        if i % 2 == 1:
            paragraph.add_run(part).bold = True
        else:
            for j, ipart in enumerate(re.split(r'\*(.+?)\*', part)):
                if not ipart:
                    continue
                run = paragraph.add_run(ipart)
                if j % 2 == 1:
                    run.italic = True

def _md_to_docx(doc, content: str):
    """Convert markdown text into python-docx paragraphs with tables and image support."""
    lines = content.strip().split("\n")
    i = 0
    while i < len(lines):
        s = lines[i].strip()

        # Markdown table detection
        if "|" in s and i + 1 < len(lines) and re.match(r"^[\|\s\-:]+$", lines[i + 1].strip()):
            # Collect table rows
            table_lines = [s]
            j = i + 2
            while j < len(lines) and "|" in lines[j]:
                table_lines.append(lines[j].strip())
                j += 1
            # Parse header
            headers = [c.strip() for c in table_lines[0].split("|") if c.strip()]
            rows = [[c.strip() for c in r.split("|") if c.strip()] for r in table_lines[1:]]
            if headers:
                from docx.shared import Pt, RGBColor
                from docx.oxml.ns import qn
                from docx.oxml import OxmlElement
                tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
                tbl.style = "Table Grid"
                hdr = tbl.rows[0]
                for ci, h in enumerate(headers):
                    cell = hdr.cells[ci]
                    cell.text = h
                    cell.paragraphs[0].runs[0].bold = True
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row[:len(headers)]):
                        tbl.rows[ri + 1].cells[ci].text = val
            i = j
            continue

        # Image embedding: ![alt](local/path)
        img_match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", s)
        if img_match:
            img_path = img_match.group(2)
            if os.path.isfile(img_path):
                try:
                    from docx.shared import Inches
                    doc.add_picture(img_path, width=Inches(5))
                except Exception:
                    doc.add_paragraph(f"[Image: {img_path}]")
            else:
                doc.add_paragraph(f"[Image not found: {img_path}]")
            i += 1
            continue

        if not s:
            doc.add_paragraph()
        elif s.startswith("#### "):
            doc.add_heading(s[5:], level=4)
        elif s.startswith("### "):
            doc.add_heading(s[4:], level=3)
        elif s.startswith("## "):
            doc.add_heading(s[3:], level=2)
        elif s.startswith("# "):
            doc.add_heading(s[2:], level=1)
        elif s.startswith("- ") or s.startswith("* "):
            _inline_fmt(doc.add_paragraph(style="List Bullet"), s[2:])
        elif re.match(r"^\d+\.\s", s):
            _inline_fmt(doc.add_paragraph(style="List Number"), re.sub(r"^\d+\.\s", "", s))
        elif s in ("---", "***", "___"):
            doc.add_page_break()
        else:
            _inline_fmt(doc.add_paragraph(), s)
        i += 1


# ── PPT helpers ───────────────────────────────────────────────────────────────

def _make_pptx(path: str, content: str):
    """Parse markdown into PowerPoint slides via python-pptx.

    Slide boundaries: '---' or start of a new '# Title' block.
    # Title     → title-only slide
    ## Subtitle → title+content slide (subtitle on title slide or section header)
    - bullet    → bullet point on current slide
    Text line   → body text on current slide
    ![alt](img) → add image to slide
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    # Layout indices: 0=title, 1=title+content, 5=blank, 6=title only
    blank_layout = prs.slide_layouts[6]   # title only
    content_layout = prs.slide_layouts[1] # title + content

    def _new_slide(title_text: str = "", is_title_slide: bool = False):
        layout = prs.slide_layouts[0] if is_title_slide else content_layout
        slide = prs.slides.add_slide(layout)
        if title_text:
            try:
                slide.shapes.title.text = title_text
            except Exception:
                pass
        return slide

    lines = content.strip().split("\n")
    current_slide = None
    current_tf = None  # text frame for body bullets
    first_slide = True

    for line in lines:
        s = line.strip()

        if s in ("---", "***"):
            current_slide = None
            current_tf = None
            continue

        if s.startswith("# "):
            title_text = s[2:].strip()
            current_slide = _new_slide(title_text, is_title_slide=first_slide)
            first_slide = False
            current_tf = None
            # Try to get content placeholder
            for ph in current_slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    current_tf = ph.text_frame
                    current_tf.clear()
                    break
            continue

        if s.startswith("## "):
            sub_text = s[3:].strip()
            if current_slide is None:
                current_slide = _new_slide(sub_text)
                first_slide = False
            else:
                # Add as subtitle / section on current slide
                for ph in current_slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        tf = ph.text_frame
                        p = tf.add_paragraph()
                        p.text = sub_text
                        p.font.bold = True
                        break
            continue

        if s.startswith("- ") or s.startswith("* "):
            bullet_text = s[2:]
            if current_slide is None:
                current_slide = _new_slide("Slide")
                first_slide = False
            if current_tf is None:
                for ph in current_slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        current_tf = ph.text_frame
                        current_tf.clear()
                        break
            if current_tf:
                p = current_tf.add_paragraph()
                p.text = bullet_text
                p.level = 0
            continue

        # Image
        img_m = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", s)
        if img_m and os.path.isfile(img_m.group(2)):
            if current_slide is None:
                current_slide = _new_slide("Image")
                first_slide = False
            try:
                current_slide.shapes.add_picture(
                    img_m.group(2), Inches(1), Inches(1.5), Inches(7)
                )
            except Exception:
                pass
            continue

        # Plain text
        if s and current_slide is not None and current_tf is not None:
            p = current_tf.add_paragraph()
            p.text = s

    if not prs.slides._sldIdLst:
        # Ensure at least one slide
        _new_slide("Presentation")

    prs.save(path)


# ── PPT conversion helper ─────────────────────────────────────────────────────

def _pptx_to_pdf(pptx_path: str, pdf_path: str) -> tuple[bool, str]:
    try:
        import win32com.client as win32
        ppt = win32.Dispatch("PowerPoint.Application")
        ppt.Visible = False
        pres = ppt.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
        pres.SaveAs(os.path.abspath(pdf_path), 32)  # 32 = ppSaveAsPDF
        pres.Close()
        ppt.Quit()
        return True, ""
    except Exception as e:
        return False, f"PPT→PDF failed: {e}. Ensure PowerPoint is installed."


# ── Excel helpers ─────────────────────────────────────────────────────────────

def _make_excel(path: str, content: str):
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = openpyxl.Workbook()
    ws = wb.active
    header_done = False
    row_num = 1
    for line in content.strip().split("\n"):
        line = line.strip()
        if not line or re.match(r"^\|?[-:]+\|", line):
            continue
        cells = ([c.strip() for c in line.split("|") if c.strip()]
                 if line.startswith("|") else
                 [c.strip() for c in line.split(",")])
        for col, val in enumerate(cells, 1):
            cell = ws.cell(row=row_num, column=col, value=val)
            if not header_done:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="2E75B6", end_color="2E75B6", fill_type="solid")
                cell.alignment = Alignment(horizontal="center")
        if not header_done:
            header_done = True
        row_num += 1
    for col in ws.columns:
        w = max((len(str(c.value or "")) for c in col), default=8)
        ws.column_dimensions[col[0].column_letter].width = min(w + 4, 50)
    wb.save(path)


# ── Conversion helpers ────────────────────────────────────────────────────────

def _docx_to_pdf(docx_path: str, pdf_path: str) -> tuple[bool, str]:
    try:
        from docx2pdf import convert
        convert(docx_path, pdf_path)
        return True, ""
    except ImportError:
        pass
    try:
        import win32com.client as win32
        word = win32.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(docx_path))
        doc.SaveAs(os.path.abspath(pdf_path), FileFormat=17)
        doc.Close(); word.Quit()
        return True, ""
    except Exception as e:
        return False, f"Conversion failed: {e}. Ensure MS Word is installed."

def _pdf_to_docx(pdf_path: str, docx_path: str) -> tuple[bool, str]:
    try:
        from pdf2docx import Converter
        cv = Converter(pdf_path)
        cv.convert(docx_path)
        cv.close()
        return True, ""
    except Exception as e:
        return False, str(e)

def _pdf_to_txt(pdf_path: str, txt_path: str) -> tuple[bool, str]:
    try:
        import fitz
        text = "".join(page.get_text() + "\n" for page in fitz.open(pdf_path))
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(text)
        return True, ""
    except Exception as e:
        return False, str(e)

def _xlsx_to_csv(xlsx_path: str, csv_path: str) -> tuple[bool, str]:
    try:
        import openpyxl, csv
        wb = openpyxl.load_workbook(xlsx_path)
        ws = wb.active
        with open(csv_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.writerow([v if v is not None else "" for v in row])
        return True, ""
    except Exception as e:
        return False, str(e)

def _csv_to_xlsx(csv_path: str, xlsx_path: str) -> tuple[bool, str]:
    try:
        import csv, openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        wb = openpyxl.Workbook()
        ws = wb.active
        with open(csv_path, encoding="utf-8") as f:
            for r, row in enumerate(csv.reader(f), 1):
                for c, val in enumerate(row, 1):
                    cell = ws.cell(row=r, column=c, value=val)
                    if r == 1:
                        cell.font = Font(bold=True, color="FFFFFF")
                        cell.fill = PatternFill(start_color="2E75B6", end_color="2E75B6", fill_type="solid")
                        cell.alignment = Alignment(horizontal="center")
        wb.save(xlsx_path)
        return True, ""
    except Exception as e:
        return False, str(e)

def _txt_to_docx(txt_path: str, docx_path: str) -> tuple[bool, str]:
    try:
        from docx import Document
        with open(txt_path, encoding="utf-8", errors="replace") as f:
            content = f.read()
        doc = Document()
        _md_to_docx(doc, content)
        doc.save(docx_path)
        return True, ""
    except Exception as e:
        return False, str(e)


# ── MCP Tools ─────────────────────────────────────────────────────────────────

@mcp.tool()
def create_file(
    filename: str,
    directory: str,
    content: str,
    file_type: str = "auto",
    confirmed: bool = False,
) -> str:
    """
    Creates a file with the given content. Supports:
      Rich documents : Word (.docx) with tables, images, bold/italic, headings
      Spreadsheets   : Excel (.xlsx), CSV
      Presentations  : PowerPoint (.pptx) - markdown to slides
      Plain-text     : HTML, JSON, JS, JSX, TS, TSX, CSS, Python,
                       YAML, SQL, XML, Markdown, TXT, and any text format
      Folder         : file_type="folder" - creates a new directory

    HARD STOP CONFIRMATION GATE (main session only):
      1. Call with confirmed=False -> preview shown -> ask user "Agreed?"
      2. END YOUR TURN. Wait for user's NEXT message.
      3. Only call with confirmed=True after user says yes.

    Parameters:
      filename  : name of the file (with or without extension, e.g. "Slides" or "app.js")
      directory : target folder - natural language OK ("Downloads", "Desktop")
                  OR absolute path
      content   : file content. Formatting:
                  Word: markdown (# H1, ## H2, - bullet, | table |, ![img](path))
                  PPT : # = slide title, ## = section, - = bullet, --- = new slide
                  Excel/CSV: comma-separated rows or markdown table
                  All others: raw text / code
      file_type : "word", "excel", "pptx", "csv", "txt", "html", "json", "python",
                  "js", "ts", "css", "yaml", "markdown", "folder", or "auto"
      confirmed : False = preview only. True = execute creation.
    """
    # Folder creation
    if file_type.lower() in ("folder", "directory"):
        dir_path = _resolve_dir(directory) or (directory if __import__("os").path.isdir(directory) else None)
        if not dir_path:
            return f"Could not locate parent directory '{directory}'."
        import os as _os
        new_folder = _os.path.join(dir_path, filename)
        if not confirmed:
            return (
                f"[CONFIRM] Create folder:\n"
                f"  \U0001f4c1 {new_folder}\n\n"
                f"Agreed?"
            )
        try:
            _os.makedirs(new_folder, exist_ok=True)
            return f"\u2713 Created folder: {new_folder}"
        except Exception as e:
            return f"Failed to create folder: {e}"

    # Resolve directory
    dir_path = _resolve_dir(directory)
    if dir_path is None:
        return (f"Could not locate directory '{directory}'. "
                "Please provide an absolute path or a more specific name.")

    ftype = _infer_type(filename, file_type)
    fname = _ensure_ext(filename, ftype)
    full_path = os.path.join(dir_path, fname)

    if not confirmed:
        ext_label = os.path.splitext(fname)[1].upper() or ftype.upper()
        return (
            f"[CONFIRM] Ready to create '{fname}' ({ext_label}) inside:\n"
            f"  \U0001f4c1 {dir_path}\n\n"
            f"Preview of content (first 300 chars):\n"
            f"{content[:300]}{'...' if len(content) > 300 else ''}\n\n"
            f"Agreed?"
        )

    # Create
    try:
        if ftype == "word":
            from docx import Document
            doc = Document()
            _md_to_docx(doc, content)
            doc.save(full_path)
        elif ftype in ("excel", "xlsx"):
            _make_excel(full_path, content)
        elif ftype in ("pptx", "ppt", "powerpoint"):
            _make_pptx(full_path, content)
        elif ftype == "csv":
            with open(full_path, "w", encoding="utf-8", newline="") as f:
                f.write(content)
        else:  # html, json, js, ts, css, yaml, py, md, txt, and all other text types
            with open(full_path, "w", encoding="utf-8") as f:
                f.write(content)

        return f"\u2713 Created '{fname}' at: {full_path}"
    except Exception as e:
        return f"Failed to create '{fname}': {e}"


@mcp.tool()
def edit_file(
    file_path: str,
    edit_instructions: str,
    append: bool = True,
    section_title: str = "",
    confirmed: bool = False,
) -> str:
    """
    Edits an existing file by appending or replacing content.

    WORKFLOW: Call with confirmed=False first → show preview → ask "Agreed?" → call with confirmed=True.

    Parameters:
      file_path        : absolute path to the file, OR a description like "AyushResume in Ayush files"
      edit_instructions: the new content to write/append (markdown for Word files)
      append           : True = add to end of file (default). False = REPLACE entire content.
      section_title    : optional heading to insert before new content (Word only)
      confirmed        : False = preview. True = execute.

    Works with: .docx, .txt, .csv, .py, .md, .log, .xlsx (appends new sheet row data)
    """
    # Resolve path if it's not absolute
    if not os.path.isfile(file_path):
        # Try to find the file using keyword
        from automation.document_parser import find_file
        user_home = os.path.expanduser("~")
        search_dirs = [
            "C:\\Ayush files",
            os.path.join(user_home, "Documents"),
            os.path.join(user_home, "Downloads"),
            os.path.join(user_home, "Desktop"),
        ]
        found = find_file(file_path, search_dirs=[d for d in search_dirs if os.path.exists(d)])
        if not found:
            return f"File not found: '{file_path}'. Provide the absolute path or check the filename."
        file_path = found

    fname = os.path.basename(file_path)
    action = "appended to" if append else "REPLACED in"

    if not confirmed:
        return (
            f"[CONFIRM] The following content will be {action} '{fname}':\n"
            f"  📄 {file_path}\n\n"
            f"Content preview (first 300 chars):\n"
            f"{edit_instructions[:300]}{'...' if len(edit_instructions) > 300 else ''}\n\n"
            f"Agreed?"
        )

    try:
        ext = os.path.splitext(file_path.lower())[1]
        if ext in (".docx", ".doc"):
            from docx import Document
            doc = Document(file_path)
            if not append:
                # Clear all paragraphs and rebuild
                for para in doc.paragraphs:
                    p = para._element
                    p.getparent().remove(p)
            if section_title:
                doc.add_heading(section_title, level=2)
            _md_to_docx(doc, edit_instructions)
            doc.save(file_path)
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            ws = wb.active
            if not append:
                ws.delete_rows(1, ws.max_row)
            for line in edit_instructions.strip().split("\n"):
                if line.strip():
                    cells = [c.strip() for c in line.split(",")]
                    ws.append(cells)
            wb.save(file_path)
        else:
            mode = "a" if append else "w"
            with open(file_path, mode, encoding="utf-8") as f:
                if append:
                    f.write("\n")
                f.write(edit_instructions)

        return f"✓ '{fname}' updated successfully: {file_path}"
    except Exception as e:
        return f"Failed to edit '{fname}': {e}"


@mcp.tool()
def delete_file(
    file_path: str,
    confirmed: bool = False,
) -> str:
    """
    Deletes a file from the filesystem.

    WORKFLOW: Call with confirmed=False → show what will be deleted → ask "Agreed?" → call with confirmed=True.

    Parameters:
      file_path : absolute path to the file, or a keyword description (e.g. "old invoice")
      confirmed : False = preview only. True = permanently delete.

    WARNING: This is irreversible. Always confirm with the user before deleting.
    """
    if not os.path.isfile(file_path):
        from automation.document_parser import find_file
        user_home = os.path.expanduser("~")
        search_dirs = [
            "C:\\Ayush files",
            os.path.join(user_home, "Documents"),
            os.path.join(user_home, "Downloads"),
            os.path.join(user_home, "Desktop"),
        ]
        found = find_file(file_path, search_dirs=[d for d in search_dirs if os.path.exists(d)])
        if not found:
            return f"File not found: '{file_path}'."
        file_path = found

    fname = os.path.basename(file_path)
    size = os.path.getsize(file_path)
    size_str = f"{size // 1024}KB" if size > 1024 else f"{size}B"

    if not confirmed:
        return (
            f"[CONFIRM] Permanently delete '{fname}' ({size_str})?\n"
            f"  📄 {file_path}\n\n"
            f"⚠️  This cannot be undone. Agreed?"
        )

    try:
        os.remove(file_path)
        return f"✓ Deleted '{fname}' from {os.path.dirname(file_path)}"
    except Exception as e:
        return f"Failed to delete '{fname}': {e}"


@mcp.tool()
def convert_file(
    file_path: str,
    target_format: str,
    output_directory: str = "",
    confirmed: bool = False,
) -> str:
    """
    Converts a file to a different format.

    WORKFLOW: Call with confirmed=False → show what will be converted → ask "Agreed?" → call with confirmed=True.

    Supported conversions:
      - Word (.docx)  → PDF
      - PDF           → Word (.docx)
      - PDF           → TXT
      - Excel (.xlsx) → CSV
      - CSV           → Excel (.xlsx)
      - TXT           → Word (.docx)

    Parameters:
      file_path        : absolute path to source file, or description/keyword
      target_format    : "pdf", "word", "docx", "txt", "excel", "xlsx", "csv"
      output_directory : where to save converted file (default: same folder as source)
      confirmed        : False = preview. True = execute.
    """
    # Resolve source file
    if not os.path.isfile(file_path):
        from automation.document_parser import find_file
        user_home = os.path.expanduser("~")
        search_dirs = [
            "C:\\Ayush files",
            os.path.join(user_home, "Documents"),
            os.path.join(user_home, "Downloads"),
            os.path.join(user_home, "Desktop"),
        ]
        found = find_file(file_path, search_dirs=[d for d in search_dirs if os.path.exists(d)])
        if not found:
            return f"Source file not found: '{file_path}'."
        file_path = found

    src_name = os.path.basename(file_path)
    src_ext = os.path.splitext(file_path.lower())[1]
    fmt = target_format.lower().strip(".")
    if fmt in ("word", "doc"):
        fmt = "docx"
    if fmt in ("ppt", "powerpoint"):
        fmt = "pptx"

    # Output path
    out_dir = output_directory if output_directory and os.path.isdir(output_directory) \
              else os.path.dirname(file_path)
    out_name = os.path.splitext(src_name)[0] + "." + fmt
    out_path = os.path.join(out_dir, out_name)

    # Validate conversion
    conv_map = {
        (".docx", "pdf"): "Word -> PDF",
        (".doc",  "pdf"): "Word -> PDF",
        (".pdf",  "docx"): "PDF -> Word",
        (".pdf",  "txt"):  "PDF -> TXT",
        (".xlsx", "csv"):  "Excel -> CSV",
        (".xls",  "csv"):  "Excel -> CSV",
        (".csv",  "xlsx"): "CSV -> Excel",
        (".txt",  "docx"): "TXT -> Word",
        (".md",   "docx"): "Markdown -> Word",
        (".pptx", "pdf"): "PowerPoint -> PDF",
        (".ppt",  "pdf"): "PowerPoint -> PDF",
    }
    conv_label = conv_map.get((src_ext, fmt))
    if not conv_label:
        return (f"Unsupported conversion: {src_ext} -> .{fmt}\n"
                f"Supported: Word->PDF, PDF->Word, PDF->TXT, Excel->CSV, CSV->Excel, TXT->Word, PPT->PDF")

    if not confirmed:
        return (
            f"[CONFIRM] {conv_label} conversion:\n"
            f"  📄 Source: {file_path}\n"
            f"  📄 Output: {out_path}\n\n"
            f"Agreed?"
        )

    # Execute conversion
    dispatch = {
        (".docx", "pdf"): lambda: _docx_to_pdf(file_path, out_path),
        (".doc",  "pdf"): lambda: _docx_to_pdf(file_path, out_path),
        (".pdf",  "docx"): lambda: _pdf_to_docx(file_path, out_path),
        (".pdf",  "txt"):  lambda: _pdf_to_txt(file_path, out_path),
        (".xlsx", "csv"):  lambda: _xlsx_to_csv(file_path, out_path),
        (".xls",  "csv"):  lambda: _xlsx_to_csv(file_path, out_path),
        (".csv",  "xlsx"): lambda: _csv_to_xlsx(file_path, out_path),
        (".txt",  "docx"): lambda: _txt_to_docx(file_path, out_path),
        (".md",   "docx"): lambda: _txt_to_docx(file_path, out_path),
        (".pptx", "pdf"): lambda: _pptx_to_pdf(file_path, out_path),
        (".ppt",  "pdf"): lambda: _pptx_to_pdf(file_path, out_path),
    }
    ok, err = dispatch[(src_ext, fmt)]()
    if ok:
        return f"✓ Converted '{src_name}' → '{out_name}'\n  Saved at: {out_path}"
    return f"Conversion failed: {err}"


@mcp.tool()
def list_directory(directory: str) -> str:
    """
    Lists all files and folders in a directory.
    directory: absolute path OR natural language (e.g. "Ayush files", "Documents", "Downloads")
    """
    dir_path = _resolve_dir(directory) or (directory if os.path.isdir(directory) else None)
    if not dir_path:
        return f"Directory not found: '{directory}'"

    try:
        entries = list(os.scandir(dir_path))
    except Exception as e:
        return f"Cannot read directory: {e}"

    folders = sorted([e for e in entries if e.is_dir()], key=lambda e: e.name.lower())
    files = sorted([e for e in entries if e.is_file()], key=lambda e: e.name.lower())
    lines = [f"📁 {dir_path}\n"]
    if folders:
        lines.append(f"Folders ({len(folders)}):")
        for e in folders:
            lines.append(f"  📁 {e.name}/")
    if files:
        lines.append(f"\nFiles ({len(files)}):")
        for e in files:
            try:
                sz = e.stat().st_size
                sz_s = f"{sz // 1024}KB" if sz > 1024 else f"{sz}B"
            except Exception:
                sz_s = "?"
            lines.append(f"  📄 {e.name} ({sz_s})")
    lines.append(f"\nTotal: {len(folders)} folders, {len(files)} files")
    return "\n".join(lines)


if __name__ == "__main__":
    mcp.run()
