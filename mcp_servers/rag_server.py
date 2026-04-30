"""
MCP Server: VoxKage Local Knowledge RAG (voxkage-rag)

Deep document memory powered by ChromaDB + sentence-transformers.
Indexes PDFs, Word docs, Excel, PowerPoint, TXT, and code files.
Detects file changes via SHA256 hash and auto-reindexes on demand.

Tools:
  index_document(file_path)          — Index or re-index a file into the RAG
  query_rag(query, top_k, file_filter) — Semantic search across all indexed documents
  check_and_index(file_path)         — Auto-index if new/changed, skip if unchanged
  list_indexed_documents()           — Show all indexed files and their metadata
  delete_from_rag(file_path)         — Remove a file from the index
  index_directory(directory, extensions) — Bulk index an entire folder/codebase

Run standalone: python mcp_servers/rag_server.py
"""

import os
import sys
import hashlib
import json
import re
import time
import logging
from pathlib import Path
from datetime import datetime

_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)

from dotenv import load_dotenv
load_dotenv(os.path.join(_ROOT, ".env"))

from mcp.server.fastmcp import FastMCP
mcp = FastMCP("voxkage-rag")
logger = logging.getLogger(__name__)

# ── Storage Paths ─────────────────────────────────────────────────────────────
_RAG_DIR   = r"C:\VoxKage\RAG"
_META_FILE = os.path.join(_RAG_DIR, "index_meta.json")
os.makedirs(_RAG_DIR, exist_ok=True)

# ── Supported file extensions ─────────────────────────────────────────────────
_TEXT_EXTS = {
    ".txt", ".md", ".csv", ".log", ".yaml", ".yml", ".json",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".css", ".html", ".htm",
    ".java", ".cpp", ".c", ".h", ".cs", ".go", ".rs", ".rb",
    ".php", ".sh", ".bat", ".ps1", ".sql", ".xml", ".toml", ".ini",
    ".gitignore", ".env",
}
_RICH_EXTS = {".pdf", ".docx", ".xlsx", ".pptx", ".xls", ".doc"}
_ALL_EXTS  = _TEXT_EXTS | _RICH_EXTS

# Skip these in bulk directory indexing
_SKIP_DIRS = {
    "__pycache__", ".git", "node_modules", ".venv", "venv",
    ".env", "dist", "build", ".next", "target", ".idea", ".vscode",
}


# ═════════════════════════════════════════════════════════════════════════════
# INTERNAL HELPERS
# ═════════════════════════════════════════════════════════════════════════════

# ── Module-level singletons (lazy-initialized, shared across all tool calls) ───
_chroma_client = None
_chroma_collection = None


def _get_collection():
    """Return the ChromaDB collection (lazy-initialized, module-level singleton)."""
    global _chroma_client, _chroma_collection
    if _chroma_collection is not None:
        return _chroma_collection
    import chromadb
    from chromadb.config import Settings
    _chroma_client = chromadb.PersistentClient(
        path=_RAG_DIR,
        settings=Settings(anonymized_telemetry=False),
    )
    _chroma_collection = _chroma_client.get_or_create_collection(
        name="voxkage_knowledge",
        metadata={"hnsw:space": "cosine"},
    )
    return _chroma_collection


def _get_embedder():
    """Return a sentence-transformer model (lazy-initialized, cached)."""
    if not hasattr(_get_embedder, "_model"):
        from sentence_transformers import SentenceTransformer
        # "all-MiniLM-L6-v2" — fast, 384-dim, good quality for local RAG
        # Downloads ~90MB on first run, then cached locally forever
        _get_embedder._model = SentenceTransformer("all-MiniLM-L6-v2")
    return _get_embedder._model


def _sha256(file_path: str) -> str:
    """Compute SHA256 hash of a file's content."""
    h = hashlib.sha256()
    try:
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return ""


def _load_meta() -> dict:
    """Load the metadata registry {file_path -> {hash, indexed_at, chunks}}."""
    try:
        if os.path.exists(_META_FILE):
            with open(_META_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}


def _save_meta(meta: dict):
    """Persist the metadata registry."""
    try:
        with open(_META_FILE, "w", encoding="utf-8") as f:
            json.dump(meta, f, indent=2, ensure_ascii=False)
    except Exception as e:
        logger.warning(f"Failed to save RAG meta: {e}")


def _extract_text(file_path: str) -> str:
    """Extract full text from a file based on its extension."""
    ext = Path(file_path).suffix.lower()

    try:
        if ext in _TEXT_EXTS:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        elif ext == ".pdf":
            import fitz
            text = ""
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text() + "\n"
            return text.strip()

        elif ext in (".docx", ".doc"):
            from docx import Document
            doc = Document(file_path)
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)
            return "\n".join(parts)

        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                parts = []
                for sheet in wb.worksheets:
                    parts.append(f"[Sheet: {sheet.title}]")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join(
                            str(c) for c in row if c is not None and str(c).strip()
                        )
                        if row_text:
                            parts.append(row_text)
                return "\n".join(parts)
            except Exception:
                return f"Could not extract Excel content from {file_path}"

        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"[Slide {i}]")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text.strip())
                return "\n".join(parts)
            except Exception:
                return f"Could not extract PowerPoint content from {file_path}"

        else:
            return f"Unsupported file type: {ext}"

    except Exception as e:
        return f"Error extracting text: {e}"


def _smart_chunk(text: str, file_path: str, chunk_size: int = 800, overlap: int = 150) -> list[str]:
    """
    Split text into overlapping chunks.
    For code files: split on function/class boundaries when possible.
    For prose: split on paragraph boundaries.
    """
    ext = Path(file_path).suffix.lower()

    if not text or not text.strip():
        return []

    # For code files: try to chunk at function/class definitions
    if ext in {".py", ".js", ".ts", ".java", ".cpp", ".c", ".cs", ".go", ".rs", ".rb", ".php"}:
        # Split on top-level function/class definitions
        pattern = r"(?m)(?=^(?:def |class |function |async def |export function |func |fn ))"
        logical_blocks = re.split(pattern, text)
        chunks = []
        current = ""
        for block in logical_blocks:
            if len(current) + len(block) > chunk_size:
                if current.strip():
                    chunks.append(current.strip())
                current = block
            else:
                current += block
        if current.strip():
            chunks.append(current.strip())
        if chunks:
            return chunks

    # For prose/documents: split on double newlines (paragraphs)
    paragraphs = re.split(r"\n{2,}", text)
    chunks = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) > chunk_size:
            if current.strip():
                chunks.append(current.strip())
            # Start new chunk with overlap from previous chunk tail
            current = current[-overlap:] + "\n\n" + para if current else para
        else:
            current += ("\n\n" if current else "") + para
    if current.strip():
        chunks.append(current.strip())

    # Final fallback: hard split
    if not chunks:
        chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size - overlap)]

    return [c for c in chunks if len(c.strip()) > 30]


def _safe_id(file_path: str, chunk_idx: int) -> str:
    """Generate a stable, ChromaDB-safe document ID."""
    key = f"{file_path}::chunk_{chunk_idx}"
    return hashlib.md5(key.encode()).hexdigest()


def _do_index(file_path: str, force: bool = False) -> dict:
    """
    Core indexing logic. Returns a status dict.
    If force=True, always re-indexes even if hash matches.
    """
    file_path = os.path.normpath(os.path.abspath(file_path))

    if not os.path.exists(file_path):
        return {"status": "error", "message": f"File not found: {file_path}"}

    ext = Path(file_path).suffix.lower()
    if ext not in _ALL_EXTS:
        return {"status": "skipped", "message": f"Unsupported extension: {ext}"}

    # Check hash for change detection
    current_hash = _sha256(file_path)
    meta = _load_meta()
    existing = meta.get(file_path, {})

    if not force and existing.get("hash") == current_hash:
        return {
            "status": "unchanged",
            "message": f"File unchanged since {existing.get('indexed_at', 'unknown')}. Using cached index.",
            "chunks": existing.get("chunks", 0),
            "file": os.path.basename(file_path),
        }

    # Extract text
    text = _extract_text(file_path)
    if not text or "Error extracting" in text or "Unsupported file" in text:
        return {"status": "error", "message": text}

    # Chunk the content
    chunks = _smart_chunk(text, file_path)
    if not chunks:
        return {"status": "error", "message": "No content extracted from file"}

    # Build embeddings
    embedder = _get_embedder()
    embeddings = embedder.encode(chunks, show_progress_bar=False).tolist()

    # Get collection and remove old chunks for this file
    collection = _get_collection()
    old_chunk_count = existing.get("chunks", 0)
    if old_chunk_count > 0:
        old_ids = [_safe_id(file_path, i) for i in range(old_chunk_count + 10)]
        try:
            collection.delete(ids=old_ids)
        except Exception:
            pass

    # Build metadata for each chunk
    file_name = os.path.basename(file_path)
    file_size = os.path.getsize(file_path)
    now = datetime.now().isoformat()

    chunk_ids = []
    metadatas = []
    for i, chunk in enumerate(chunks):
        cid = _safe_id(file_path, i)
        chunk_ids.append(cid)
        metadatas.append({
            "file_path": file_path,
            "file_name": file_name,
            "file_ext": ext,
            "chunk_index": i,
            "total_chunks": len(chunks),
            "indexed_at": now,
            "file_hash": current_hash,
            "file_size_kb": round(file_size / 1024, 1),
        })

    # Upsert into ChromaDB
    collection.upsert(
        ids=chunk_ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=metadatas,
    )

    # Update metadata registry
    action = "reindexed" if existing else "indexed"
    meta[file_path] = {
        "hash": current_hash,
        "indexed_at": now,
        "chunks": len(chunks),
        "file_name": file_name,
        "file_ext": ext,
        "file_size_kb": round(file_size / 1024, 1),
    }
    _save_meta(meta)

    return {
        "status": "success",
        "action": action,
        "file": file_name,
        "chunks": len(chunks),
        "characters": len(text),
        "message": f"[OK] {action.capitalize()} '{file_name}' into RAG ({len(chunks)} chunks from {len(text):,} characters)",
    }


# ═════════════════════════════════════════════════════════════════════════════
# MCP TOOLS
# ═════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def index_document(file_path: str) -> str:
    """
    RAG MEMORY: Index or re-index a document into VoxKage's local knowledge base.

    Automatically detects if the file is NEW (indexes fresh) or CHANGED (reindexes)
    or UNCHANGED (returns cached result instantly — no work done).

    Supports: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), TXT, CSV,
              Python, JavaScript, TypeScript, Java, C/C++, Go, Rust, SQL, Markdown, YAML, JSON, and more.

    ALWAYS call this after reading any file for the first time so VoxKage remembers it.
    """
    result = _do_index(file_path, force=False)
    return result.get("message", str(result))


@mcp.tool()
def check_and_index(file_path: str) -> str:
    """
    RAG MEMORY: Smart auto-index gate. Call this BEFORE reading ANY file.

    - If file is new → indexes it automatically
    - If file has changed since last index → reindexes it
    - If file is unchanged → returns instantly from cache

    This is the STANDARD workflow for file access:
    1. check_and_index(file_path)   ← always call this first
    2. query_rag(your_question)     ← search the indexed content
    """
    result = _do_index(file_path, force=False)
    status = result.get("status", "error")
    if status == "unchanged":
        return f"[RAG] Already indexed and unchanged — using cache. {result['chunks']} chunks available for '{result['file']}'."
    elif status == "success":
        return f"[RAG] {result['message']}"
    elif status == "skipped":
        return f"[RAG SKIP] {result['message']} — will read file directly."
    else:
        return f"[RAG ERROR] {result.get('message', 'Unknown error')}"


@mcp.tool()
def query_rag(
    query: str,
    top_k: int = 5,
    file_filter: str = "",
) -> str:
    """
    RAG MEMORY: Semantic search across all indexed documents.

    Use this to answer questions about ANY file VoxKage has seen before,
    without needing to re-read the file.

    Parameters:
      query       : Natural language question or keyword search
      top_k       : Number of results to return (default 5)
      file_filter : Optional — restrict search to files matching this name/path substring

    Examples:
      query_rag("what is the tax rate in the invoices")
      query_rag("how did we handle API fallbacks", file_filter="web_agent")
      query_rag("find the authentication section", file_filter="auth")
    """
    try:
        collection = _get_collection()
        count = collection.count()
        if count == 0:
            return "[RAG] No documents indexed yet. Use index_document() or check_and_index() first."

        embedder = _get_embedder()
        query_embedding = embedder.encode([query], show_progress_bar=False).tolist()[0]

        where_filter = None
        if file_filter:
            # ChromaDB 'where' on metadata — use $contains logic via Python post-filter
            pass

        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=min(top_k * 2, count),  # over-fetch if filtering
            include=["documents", "metadatas", "distances"],
        )

        docs = results.get("documents", [[]])[0]
        metas = results.get("metadatas", [[]])[0]
        distances = results.get("distances", [[]])[0]

        if not docs:
            return "[RAG] No relevant results found for that query."

        # Apply file_filter post-query if given
        filtered = []
        for doc, meta, dist in zip(docs, metas, distances):
            if file_filter:
                fp = meta.get("file_path", "").lower()
                fn = meta.get("file_name", "").lower()
                if file_filter.lower() not in fp and file_filter.lower() not in fn:
                    continue
            filtered.append((doc, meta, dist))
            if len(filtered) >= top_k:
                break

        if not filtered:
            return f"[RAG] No results matching file_filter='{file_filter}'. Try without the filter."

        lines = [f"[RAG] Found {len(filtered)} relevant chunk(s) for: \"{query}\"\n"]
        for i, (doc, meta, dist) in enumerate(filtered, 1):
            confidence = round((1 - dist) * 100, 1)
            fn = meta.get("file_name", "unknown")
            chunk_n = meta.get("chunk_index", 0)
            total = meta.get("total_chunks", 1)
            lines.append(f"--- Result {i} | {fn} (chunk {chunk_n + 1}/{total}) | {confidence}% relevance ---")
            lines.append(doc.strip())
            lines.append("")

        return "\n".join(lines)

    except Exception as e:
        return f"[RAG ERROR] Query failed: {e}"


@mcp.tool()
def list_indexed_documents() -> str:
    """
    RAG MEMORY: List all documents currently in VoxKage's knowledge base.

    Shows: file name, size, chunk count, when indexed, and whether the file has
    changed on disk since it was last indexed (needs reindex if so).
    """
    meta = _load_meta()
    if not meta:
        return "[RAG] No documents indexed yet. Use index_document() to start building your knowledge base."

    lines = [f"[RAG] VoxKage Knowledge Base — {len(meta)} indexed document(s)\n"]
    lines.append(f"{'File':<45} {'Chunks':>6}  {'Size KB':>8}  {'Indexed At':<22}  {'Status'}")
    lines.append("-" * 110)

    for file_path, info in sorted(meta.items(), key=lambda x: x[1].get("indexed_at", ""), reverse=True):
        fn = info.get("file_name", os.path.basename(file_path))[:44]
        chunks = info.get("chunks", "?")
        size_kb = info.get("file_size_kb", "?")
        indexed_at = info.get("indexed_at", "unknown")[:19]

        # Check if file has changed
        status = "[OK] current"
        if not os.path.exists(file_path):
            status = "[!!] FILE MISSING"
        elif _sha256(file_path) != info.get("hash", ""):
            status = "[!!] CHANGED -- needs reindex"

        lines.append(f"{fn:<45} {str(chunks):>6}  {str(size_kb):>8}  {indexed_at:<22}  {status}")

    return "\n".join(lines)


@mcp.tool()
def delete_from_rag(file_path: str) -> str:
    """
    RAG MEMORY: Remove a document from VoxKage's knowledge base.
    Use this if a file has been deleted or you no longer want it indexed.
    """
    file_path = os.path.normpath(os.path.abspath(file_path))
    meta = _load_meta()

    if file_path not in meta:
        return f"[RAG] '{os.path.basename(file_path)}' is not in the index."

    info = meta[file_path]
    chunk_count = info.get("chunks", 0)

    # Delete chunks from ChromaDB
    try:
        collection = _get_collection()
        ids_to_delete = [_safe_id(file_path, i) for i in range(chunk_count + 10)]
        collection.delete(ids=ids_to_delete)
    except Exception as e:
        logger.warning(f"ChromaDB delete error: {e}")

    # Remove from metadata
    del meta[file_path]
    _save_meta(meta)

    return f"[RAG] Removed '{info.get('file_name', os.path.basename(file_path))}' and its {chunk_count} chunks from the knowledge base."


@mcp.tool()
def index_directory(
    directory: str,
    extensions: str = "",
    recursive: bool = True,
) -> str:
    """
    RAG MEMORY: Bulk-index an entire directory (codebase, documents folder, etc.).

    Parameters:
      directory  : The folder to index (e.g. "C:\\Users\\AYUSH\\Desktop\\Vision-Assistant")
      extensions : Comma-separated list of file extensions to include (e.g. ".py,.md,.txt")
                   Leave empty to use all supported extensions.
      recursive  : Whether to index subdirectories (default True)

    This is how you index an ENTIRE CODEBASE or documents folder.
    Already-indexed unchanged files are skipped automatically (only new/changed files are processed).
    """
    directory = os.path.normpath(os.path.abspath(directory))
    if not os.path.isdir(directory):
        return f"[RAG ERROR] Not a directory: {directory}"

    # Parse extension filter
    if extensions:
        target_exts = {e.strip().lower() if e.strip().startswith(".") else f".{e.strip().lower()}"
                       for e in extensions.split(",")}
        target_exts = target_exts & _ALL_EXTS
    else:
        target_exts = _ALL_EXTS

    results = {"indexed": 0, "reindexed": 0, "unchanged": 0, "skipped": 0, "errors": 0}
    processed = []

    walker = os.walk(directory) if recursive else [(directory, [], os.listdir(directory))]

    for root, dirs, files in walker:
        # Prune skip directories in-place
        if recursive:
            dirs[:] = [d for d in dirs if d not in _SKIP_DIRS and not d.startswith(".")]

        for fname in files:
            ext = Path(fname).suffix.lower()
            if ext not in target_exts:
                continue

            fpath = os.path.join(root, fname)
            try:
                result = _do_index(fpath, force=False)
                status = result.get("status", "error")
                action = result.get("action", "")
                if status == "success" and action == "indexed":
                    results["indexed"] += 1
                    processed.append(f"  [NEW]: {fname}")
                elif status == "success" and action == "reindexed":
                    results["reindexed"] += 1
                    processed.append(f"  [UPDATED]: {fname}")
                elif status == "unchanged":
                    results["unchanged"] += 1
                elif status == "skipped":
                    results["skipped"] += 1
                else:
                    results["errors"] += 1
                    processed.append(f"  [ERROR]: {fname} -- {result.get('message', '')}")
            except Exception as e:
                results["errors"] += 1
                processed.append(f"  [EXCEPTION]: {fname} -- {e}")

    summary_lines = [
        f"[RAG] Directory indexing complete for: {directory}",
        f"  New files indexed  : {results['indexed']}",
        f"  Files reindexed    : {results['reindexed']}",
        f"  Unchanged (skipped): {results['unchanged']}",
        f"  Unsupported/skipped: {results['skipped']}",
        f"  Errors             : {results['errors']}",
        "",
    ]
    if processed:
        summary_lines.append("Changes:")
        summary_lines.extend(processed[:30])  # cap output
        if len(processed) > 30:
            summary_lines.append(f"  ... and {len(processed) - 30} more")

    return "\n".join(summary_lines)


if __name__ == "__main__":
    mcp.run()
